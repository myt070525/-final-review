import os


def extract_text_from_file(filepath: str, ext: str) -> str:
    if ext == ".txt":
        with open(filepath, "r", encoding="utf-8") as f:
            return f.read()

    elif ext == ".docx":
        from docx import Document
        doc = Document(filepath)
        return "\n".join(p.text for p in doc.paragraphs if p.text.strip())

    elif ext == ".pptx":
        from pptx import Presentation
        prs = Presentation(filepath)
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        if para.text.strip():
                            texts.append(para.text.strip())
        return "\n".join(texts)

    elif ext == ".pdf":
        import fitz  # PyMuPDF
        doc = fitz.open(filepath)
        texts = []
        for page in doc:
            t = page.get_text()
            if t.strip():
                texts.append(t.strip())
        doc.close()
        return "\n".join(texts)

    raise ValueError(f"Unsupported format: {ext}")
